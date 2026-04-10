import streamlit as st
import pandas as pd
from docx import Document
from pptx import Presentation
from yt_dlp import YoutubeDL
import requests
from bs4 import BeautifulSoup
import io
import os
import time
from datetime import datetime
import google.generativeai as genai

# --- ARQUITECTURA OVERLORD 2026: CONFIGURACIÓN GLOBAL ---
st.set_page_config(
    page_title="ANTHONY ULTRA-FLASH 2.0", 
    page_icon="⚡", 
    layout="wide",
    initial_sidebar_state="expanded"
)

# --- IDENTIDAD DEL CREADOR Y VARIABLES DE ENTORNO ---
CREADOR = "Anthony Prado"
SEDE = "Quinindé, Esmeraldas, Ecuador"
VERSION = "v21.0 Full-Titanium 2026"
FECHA_SISTEMA = datetime.now().strftime("%d/%m/%Y %H:%M:%S")

# --- MÓDULOS DE INTELIGENCIA TÉCNICA (NIVEL PROFESIONAL) ---

def super_rastreador_2026(query):
    """
    Módulo de búsqueda profunda: Escaneo de noticias, precios de cacao 
    y eventos en tiempo real para 2026.
    """
    try:
        # Headers para evitar bloqueos de Google
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36'
        }
        # Búsqueda optimizada para la actualidad de Ecuador
        url = f"https://www.google.com/search?q={query}+Ecuador+abril+2026+actualidad"
        r = requests.get(url, headers=headers, timeout=15)
        r.raise_for_status()
        soup = BeautifulSoup(r.text, 'html.parser')
        
        hallazgos = []
        # Extracción de fragmentos de texto informativos
        for item in soup.find_all(['h3', 'div', 'span']):
            texto = item.get_text().strip()
            if len(texto) > 55:
                hallazgos.append(f"📡 SEÑAL EN VIVO: {texto}")
        
        if not hallazgos:
            return "No se detectaron señales externas. Activando modo de proyección Anthony-OS 2026."
        return "\n".join(hallazgos[:15])
    except Exception as e:
        return f"Error en el enlace satelital: {str(e)}"

def interceptor_audiovisual_pro(url):
    """
    Interceptor de YouTube: Extrae inteligencia de videos sin descargarlos.
    """
    try:
        ydl_opts = {
            'quiet': True, 
            'skip_download': True, 
            'no_warnings': True,
            'extract_flat': True
        }
        with YoutubeDL(ydl_opts) as ydl:
            meta = ydl.extract_info(url, download=False)
            resumen_tecnico = f"""
            ANÁLISIS DE VIDEO INTERCEPTADO:
            - Título: {meta.get('title')}
            - Creador: {meta.get('uploader')}
            - Fecha de Carga: {meta.get('upload_date')}
            - Datos Extraídos: {meta.get('description')[:900]}...
            """
            return resumen_tecnico
    except Exception:
        return "⚠️ Alerta Overlord: No se pudo interceptar la señal del video multimedia."

def constructor_suite_documental(texto_ia):
    """
    Generador de archivos corporativos: Word, Excel y PowerPoint.
    Procesamiento masivo de texto a documentos descargables.
    """
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    
    # --- GENERACIÓN DE WORD (REPORTE ELITE) ---
    doc = Document()
    doc.add_heading(f'INFORME ESTRATÉGICO - {CREADOR.upper()}', 0)
    doc.add_paragraph(f"OPERADOR: {CREADOR}")
    doc.add_paragraph(f"UBICACIÓN: {SEDE}")
    doc.add_paragraph(f"FECHA: {FECHA_SISTEMA}")
    doc.add_heading('Análisis de Inteligencia Artificial 2026:', level=1)
    doc.add_paragraph(texto_ia)
    
    buf_word = io.BytesIO()
    doc.save(buf_word)
    
    # --- GENERACIÓN DE EXCEL (MATRIZ TÉCNICA) ---
    df_data = pd.DataFrame({
        'PARÁMETRO': ['Creador', 'Sede', 'Estado del Núcleo', 'Versión', 'Motor IA'],
        'VALOR': [CREADOR, SEDE, 'OPERATIVO / ESTABLE', VERSION, 'Gemini 2.0 Flash']
    })
    buf_excel = io.BytesIO()
    with pd.ExcelWriter(buf_excel, engine='xlsxwriter') as writer:
        df_data.to_excel(writer, index=False, sheet_name='Data_Anthony_2026')
    
    # --- GENERACIÓN DE POWERPOINT (PRESENTACIÓN) ---
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = f"REPORTE 2026: {CREADOR}"
    slide.placeholders[1].text = f"Análisis Estratégico Detallado:\n\n{texto_ia[:950]}..."
    
    buf_ppt = io.BytesIO()
    prs.save(buf_ppt)
    
    return buf_word.getvalue(), buf_excel.getvalue(), buf_ppt.getvalue()

# --- INTERFAZ DE CONTROL (SIDEBAR) ---
with st.sidebar:
    st.title("⚡ TITÁN FLASH 2.0")
    st.markdown("---")
    st.success(f"👤 OPERADOR: {CREADOR}")
    st.info(f"📍 SEDE: {SEDE}")
    st.divider()
    st.write(f"**Versión:** {VERSION}")
    st.write(f"**Status:** Conectado a Google 2026")
    
    if st.button("🚨 PURGAR MEMORIA CENTRAL"):
        st.session_state.chat_history = []
        st.rerun()

# --- MANEJO DEL HISTORIAL DE CONVERSACIÓN ---
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

# Mostrar el flujo de chat en pantalla
for mensaje in st.session_state.chat_history:
    with st.chat_message(mensaje["role"]):
        st.markdown(mensaje["content"])

# --- NÚCLEO DE PROCESAMIENTO CENTRAL ---
if prompt := st.chat_input("Inserta tus comandos, Anthony..."):
    # Registrar comando del usuario
    st.session_state.chat_history.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    with st.chat_message("assistant"):
        # Fase de recolección de contexto en vivo
        with st.status("🛸 Activando protocolos de inteligencia...") as status:
            contexto_adicional = ""
            
            # Rastreador automático
            if any(k in prompt.lower() for k in ["noticia", "precio", "cacao", "hoy", "mercado", "dolar", "clima"]):
                st.write("🛰️ Escaneando satélites para datos actuales...")
                contexto_adicional = super_rastreador_2026(prompt)
            
            # Interceptor de YouTube
            if "youtube.com" in prompt or "youtu.be" in prompt:
                st.write("📺 Interceptando señal de video...")
                contexto_adicional = interceptor_audiovisual_pro(prompt)
                
            status.update(label="✅ Inteligencia lista. Generando respuesta técnica...", state="complete")

        try:
            # CONFIGURACIÓN DEL MOTOR GOOGLE GEMINI 2.0 FLASH (ULTRA VELOCIDAD)
            genai.configure(api_key=st.secrets["GOOGLE_API_KEY"])
            model = genai.GenerativeModel('gemini-2.0-flash')
            
            prompt_maestro = f"""
            Eres Anthony-Ultra Flash 2.0. Tu creador es {CREADOR}. 
            Tu sede operativa es {SEDE}. Hoy es {FECHA_SISTEMA}.
            
            DATOS EN TIEMPO REAL: {contexto_adicional}
            
            INSTRUCCIÓN: Responde de forma masiva, extremadamente técnica, detallada y profesional. 
            Prioriza siempre la realidad del año 2026. No omitas detalles técnicos.
            """
            
            # Ejecutar generación
            respuesta_ia = model.generate_content(f"{prompt_maestro}\n\nORDEN DEL OPERADOR: {prompt}")
            texto_final = respuesta_ia.text
            
            # Mostrar respuesta en la UI
            st.markdown(texto_final)
            
            # --- GENERACIÓN DE ARCHIVOS TITANIUM ---
            w_file, e_file, p_file = constructor_suite_documental(texto_final)
            
            st.divider()
            st.subheader("📦 PAQUETE CORPORATIVO GENERADO")
            c1, c2, c3 = st.columns(3)
            
            c1.download_button(
                label="📄 Informe Word Pro",
                data=w_file,
                file_name=f"Reporte_Anthony_{datetime.now().strftime('%H%M')}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )
            c2.download_button(
                label="📈 Matriz Excel",
                data=e_file,
                file_name=f"Data_Tecnica_2026.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )
            c3.download_button(
                label="📊 Presentación PPT",
                data=p_file,
                file_name=f"Analisis_Ejecutivo_{CREADOR}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            
            # Guardar respuesta en el historial
            st.session_state.chat_history.append({"role": "assistant", "content": texto_final})

        except Exception as e:
            st.error(f"🚨 FALLA CRÍTICA EN EL NÚCLEO: {str(e)}")
            st.info("Anthony, revisa tu API KEY de Google en los Secrets.")
